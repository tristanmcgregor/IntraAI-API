from __future__ import annotations

from pathlib import Path
from typing import Iterable, List, Tuple, Optional, TypedDict
from io import BytesIO
import os
import sqlite3
import hashlib
import re

from bs4 import BeautifulSoup
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from .config import Settings


def _clean_text(text: str) -> str:
	# Collapse multiple spaces
	text = re.sub(r"[ \t\u00A0]{2,}", " ", text)
	# Fix hyphenation at line breaks: e.g., "ware-\nhouse" -> "warehouse"
	text = re.sub(r"(\w)-\s*\n\s*(\w)", r"\1\2", text)
	# Normalize line breaks
	text = re.sub(r"\r\n|\r", "\n", text)
	# Remove repeated short fragments across line breaks
	lines = [l.strip() for l in text.split("\n") if l.strip()]
	dedup = []
	seen = set()
	for l in lines:
		key = l.lower()
		if key in seen:
			continue
		seen.add(key)
		dedup.append(l)
	return "\n".join(dedup)


def _load_txt_bytes(data: bytes) -> str:
	return data.decode("utf-8", errors="ignore")


def _load_pdf_bytes(data: bytes) -> str:
	reader = PdfReader(BytesIO(data))
	texts: List[str] = []
	for page in reader.pages:
		texts.append(page.extract_text() or "")
	return "\n".join(texts)


def _load_docx_bytes(data: bytes) -> str:
	doc = DocxDocument(BytesIO(data))
	paras = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
	return "\n".join(paras)


def _load_pptx_bytes(data: bytes) -> str:
	prs = Presentation(BytesIO(data))
	texts: List[str] = []
	for slide in prs.slides:
		for shape in slide.shapes:
			if hasattr(shape, "text") and shape.text:
				texts.append(shape.text)
	return "\n".join(texts)


def _load_html_bytes(data: bytes) -> str:
	soup = BeautifulSoup(data.decode("utf-8", errors="ignore"), "lxml")
	return soup.get_text(" ", strip=True)


class StoredDocument(TypedDict):
    source: str
    bytes: bytes
    content_type: str
    s3_key: str


def load_documents_from_records(documents: Iterable[StoredDocument]) -> List[tuple[str, dict]]:
    loaded: List[tuple[str, dict]] = []
    for doc in documents:
        source_name = doc.get("source")
        blob = doc.get("bytes")
        if not blob:
            print(f"[ingest] Missing bytes for {source_name}; skipping")
            continue
        suffix = Path(source_name).suffix.lower() if source_name else ""
        try:
            if suffix in {".txt", ".md", ".rst"}:
                text = _load_txt_bytes(blob)
            elif suffix == ".pdf":
                text = _load_pdf_bytes(blob)
            elif suffix == ".docx":
                text = _load_docx_bytes(blob)
            elif suffix == ".pptx":
                text = _load_pptx_bytes(blob)
            elif suffix in {".htm", ".html"}:
                text = _load_html_bytes(blob)
            else:
                text = _load_txt_bytes(blob)
            text = _clean_text(text)
            if text.strip():
                metadata = {
                    "source": source_name,
                    "path": doc.get("s3_key", source_name),
                }
                if "s3_key" in doc:
                    metadata["s3_key"] = doc["s3_key"]
                loaded.append((text, metadata))
        except Exception as exc:
            print(f"[ingest] Failed to parse {source_name}: {exc}")
    return loaded


def chunk_text(text: str, chunk_size: int, chunk_overlap: int) -> List[str]:
	chunks: List[str] = []
	start = 0
	while start < len(text):
		end = min(start + chunk_size, len(text))
		chunks.append(text[start:end])
		if end == len(text):
			break
		start = end - chunk_overlap
		if start < 0:
			start = 0
	return chunks


def chunk_documents(documents: List[tuple[str, dict]], settings: Settings) -> List[tuple[str, dict]]:
	result: List[tuple[str, dict]] = []
	for text, meta in documents:
		for chunk in chunk_text(text, settings.chunk_size, settings.chunk_overlap):
			result.append((chunk, meta))
	return result


def ensure_schema(conn: sqlite3.Connection) -> None:
	cur = conn.cursor()
	cur.execute(
		"""
		CREATE TABLE IF NOT EXISTS chunks (
			id TEXT PRIMARY KEY,
			source TEXT,
			path TEXT,
			content TEXT,
			embedding BLOB
		);
		"""
	)
	cur.execute("CREATE INDEX IF NOT EXISTS idx_chunks_source ON chunks(source)")
	conn.commit()


def upsert_chunk(conn: sqlite3.Connection, chunk_id: str, source: str, path: str, content: str, embedding: bytes) -> None:
	cur = conn.cursor()
	cur.execute(
		"""
		INSERT INTO chunks(id, source, path, content, embedding)
		VALUES(?,?,?,?,?)
		ON CONFLICT(id) DO UPDATE SET source=excluded.source, path=excluded.path, content=excluded.content, embedding=excluded.embedding
		""",
		(chunk_id, source, path, content, embedding),
	)
	conn.commit()


def compute_id(source: str, content: str) -> str:
	m = hashlib.sha256()
	m.update(source.encode("utf-8"))
	m.update(b"\n")
	m.update(content.encode("utf-8"))
	return m.hexdigest()


def chunk_exists(conn: sqlite3.Connection, chunk_id: str) -> bool:
	cur = conn.cursor()
	cur.execute("SELECT 1 FROM chunks WHERE id=?", (chunk_id,))
	return cur.fetchone() is not None
